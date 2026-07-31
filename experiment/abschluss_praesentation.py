"""Abschlusspräsentation — von Grund auf neues Deck.

Eigenständiges Design ohne Bezug zu früheren Decks: Action-Titles,
Abschnitts-Trennfolien mit Großtypografie, Kicker-Zeilen, Akzentbalken
und Hairlines statt Karten und Pills; Punkt-, Hantel- und Matrix-Charts
statt Balkendiagrammen; Farbwelt Violett/Teal/Magenta (CVD-validiert).

Inhalte: sechzehn kompakte Inhaltsfolien, vier kurze Abschnittstrenner,
Fragenfolie und eine Backupfolie. Die vorgegebene Vierer-Gliederung bleibt im
Hauptteil vollständig sichtbar: Projektthema, Projektverlauf, Projektstand
sowie Zusammenfassung und Ausblick. Alle Ergebnisansichten bleiben vor der
Fragenfolie; nur die ausführliche Metrikfehler-Folie folgt als Backup. Das
qualitative Video zeigt ein komplettes ungesehenes CV-Testexperiment.

Aufbau: python experiment/abschluss_praesentation.py
"""

from __future__ import annotations

import argparse
import tempfile
from pathlib import Path

import matplotlib

matplotlib.use("Agg")
import matplotlib.pyplot as plt
from matplotlib.colors import LinearSegmentedColormap
from matplotlib.lines import Line2D
from PIL import Image

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

ROOT = Path(__file__).resolve().parent.parent
TEMPLATE = ROOT / "SlideMaster_Template_16_9.pptx"
# Assets aus experiment/detektionsvideo_cv_three_views.py (vorher ausführen)
VIDEOS = ROOT / "results" / "videos"
VIDEO = VIDEOS / "exp1_cv_full_three_views_oblique.mp4"
POSTER = VIDEOS / "exp1_cv_full_three_views_oblique_poster.png"
SCENE_OS1 = VIDEOS / "exp1_scene_os1.png"
SCENE_MERGED = VIDEOS / "exp1_scene_merged.png"
DEFAULT_OUTPUT = ROOT / "Abschlusspräsentation DCAITI LiDAR-Projekt.pptx"

# ------------------------------------------------------------ Designsystem
INK = RGBColor(0x1B, 0x1E, 0x24)
SUBTLE = RGBColor(0x5F, 0x66, 0x72)
FAINT = RGBColor(0x9A, 0xA1, 0xAC)
HAIR = RGBColor(0xE4, 0xE6, 0xEA)
ACCENT = RGBColor(0x6E, 0x56, 0xCF)
ACCENT_DEEP = RGBColor(0x4B, 0x3A, 0x99)
TINT = RGBColor(0xF0, 0xED, 0xFA)
GHOST = RGBColor(0xDD, 0xD5, 0xF4)
TRACK = RGBColor(0xEF, 0xF0, 0xF3)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

# Chart-Hexwerte (validiert: dataviz-Validator, light mode, alle Checks PASS)
CX_INK = "#1B1E24"
CX_SUBTLE = "#5F6672"
CX_FAINT = "#9AA1AC"
CX_HAIR = "#D8DBE0"
CX_ACCENT = "#6E56CF"
CX_GRAY = "#A6ACB8"
V_MERGED = "#6E56CF"
V_OS0 = "#0E9888"
V_OS1 = "#9F1853"
VIEW_HEX = [V_MERGED, V_OS0, V_OS1]
VIEW_RGB = [RGBColor(0x6E, 0x56, 0xCF), RGBColor(0x0E, 0x98, 0x88),
            RGBColor(0x9F, 0x18, 0x53)]
VIEWS = ["Fusion (merged)", "os0", "os1"]

plt.rcParams.update({
    "font.family": "Segoe UI",
    "text.color": CX_INK,
    "axes.edgecolor": CX_HAIR,
    "axes.labelcolor": CX_SUBTLE,
    "xtick.color": CX_SUBTLE,
    "ytick.color": CX_INK,
})

# ------------------------------------------------------------------ Daten
# PointPillars, experiment-held-out 3-Fold-CV (results/CROSS_VALIDATION_RESULTS.md)
PP_FOLDS = {"merged": [0.840, 0.786, 0.768],
            "os0": [0.793, 0.781, 0.757],
            "os1": [0.801, 0.751, 0.759]}
PP_MEAN = {"merged": 0.798, "os0": 0.777, "os1": 0.770}
# gepoolt out-of-fold, nur bewegte Objekte: Klasse -> [(AP30, AP60)] je Sicht.
# Label-basierte Bewegt-Wertung (results/CROSS_VALIDATION_STATIC_AWARE.json):
# Boxen auf statisch gelabelten Objekten (IoU >= 0.25) zählen nicht als
# Fehlalarm — die manuellen bewegt/statisch-Labels machen sie zuordenbar.
# 4 Nachkommastellen aus dem JSON, damit die 2-stellige Anzeige korrekt
# rundet.
DYN = {
    "Person": [(0.9015, 0.7288), (0.6525, 0.3697), (0.8959, 0.4257)],
    "Fahrrad": [(0.7999, 0.6677), (0.8945, 0.3783), (0.7444, 0.5285)],
    "Auto": [(0.8910, 0.8832), (0.8853, 0.8563), (0.8316, 0.7891)],
}
# Gepoolte out-of-fold mAP nur für bewegte Objekte (AP30–60, 3 Klassen),
# gleiche label-basierte Wertung
DYN_MAP = {"merged": 0.8355, "os0": 0.7096, "os1": 0.7560}
# AP30 der bewegten Zielklasse je ungesehenem Testexperiment
EXP_AP30 = {
    "merged": [0.908, 0.898, 0.907, 0.904, 0.727, 0.804, 0.908, 0.882, 0.902],
    "os0": [0.921, 0.964, 0.942, 0.906, 0.907, 0.893, 0.908, 0.885, 0.903],
    "os1": [0.873, 0.768, 0.878, 0.815, 0.718, 0.731, 0.899, 0.902, 0.905],
}
EXP_COLS = [("1", "Auto"), ("2", "Auto"), ("3", "Auto"), ("4", "Rad"),
            ("5", "Rad"), ("6", "Rad"), ("7", "Person"), ("8", "Person"),
            ("9", "Person")]
# CenterPoint, identisches Protokoll (results/CENTERPOINT_CV_RESULTS.md)
CP_MEAN = {"merged": 0.751, "os0": 0.659, "os1": 0.723}
# Geisterboxen out-of-fold (PointPillars), Score >= 0.3 und kein noch
# ungematchtes GT derselben Klasse bei 3D-IoU >= 0.3. Exakte Zählung:
# results/GHOST_BOXES.json. Sicht -> Klasse -> Gesamtzahl.
GHOST_BOXES = {
    "merged": {"Person": 510, "Fahrrad": 315, "Auto": 159},
    "os0": {"Person": 948, "Fahrrad": 157, "Auto": 415},
    "os1": {"Person": 329, "Fahrrad": 372, "Auto": 248},
}
GHOST_CLASS_HEX = [V_MERGED, V_OS0, V_OS1]
RUNTIME_MS = {"merged": 85, "os0": 45, "os1": 43}
PRETRAIN = [("Diese Arbeit (Finetuning, CV)", 0.798, True),
            ("nur vortrainiert · KITTI", 0.058, False),
            ("nur vortrainiert · LUMPI", 0.045, False),
            ("nur vortrainiert · OSDaR23", 0.0, False)]


# ------------------------------------------------------------------ Charts

def _strip(ax, keep_x=True):
    for side in ("top", "right", "left"):
        ax.spines[side].set_visible(False)
    ax.spines["bottom"].set_visible(keep_x)
    ax.tick_params(length=0)


def fig_map_dots(path: Path) -> None:
    """Cleveland-Dot-Plot: Fold-Werte + Fold-Mittel je Sicht."""
    fig, ax = plt.subplots(figsize=(6.2, 2.2), dpi=250)
    order = ["merged", "os0", "os1"]
    for row, (view, color, name) in enumerate(zip(order, VIEW_HEX, VIEWS)):
        mean = PP_MEAN[view]
        ax.plot([0.70, mean], [row, row], color=color, linewidth=2.4,
                alpha=0.42, zorder=2, solid_capstyle="round")
        ax.scatter([mean], [row], s=105, marker="D", color=color,
                   edgecolor="white", linewidths=0.8, zorder=4)
        ax.text(mean + 0.010, row, f"{mean:.3f}", va="center",
                ha="left", fontsize=12.5, fontweight="bold", color=CX_INK)
    ax.set_yticks(range(3))
    ax.set_yticklabels(VIEWS, fontsize=12)
    ax.invert_yaxis()
    ax.set_ylim(2.55, -0.55)
    ax.set_xlim(0.70, 0.90)
    ax.set_xticks([0.70, 0.75, 0.80, 0.85, 0.90])
    ax.set_xticklabels(["0.70", "0.75", "0.80", "0.85", "0.90"], fontsize=10)
    ax.set_xlabel("Fold-Mittel mAP (AP30–60) · Achse gezoomt",
                  fontsize=9.5)
    _strip(ax)
    ax.grid(axis="x", color="#EFF0F3", linewidth=1.1)
    ax.set_axisbelow(True)
    fig.tight_layout(pad=0.4)
    fig.savefig(path, dpi=250)
    plt.close(fig)


def fig_dynamic_map(path: Path) -> None:
    """Dot-Plot: gepoolte Gesamt-mAP nur für bewegte Objekte."""
    fig, ax = plt.subplots(figsize=(6.1, 2.55), dpi=250)
    order = ["merged", "os0", "os1"]
    for row, (view, color, name) in enumerate(zip(order, VIEW_HEX, VIEWS)):
        value = DYN_MAP[view]
        ax.plot([0.65, value], [row, row], color=color, linewidth=2.4,
                alpha=0.42, zorder=2, solid_capstyle="round")
        ax.scatter([value], [row], s=125, color=color, edgecolor="white",
                   linewidths=0.9, zorder=3)
        ax.text(value + 0.008, row, f"{value:.3f}", va="center", ha="left",
                fontsize=12.5, fontweight="bold", color=CX_INK)
    ax.set_yticks(range(3))
    ax.set_yticklabels(VIEWS, fontsize=12)
    ax.set_ylim(2.55, -0.55)
    ax.set_xlim(0.65, 0.84)
    ax.set_xticks([0.65, 0.70, 0.75, 0.80])
    ax.set_xticklabels(["0.65", "0.70", "0.75", "0.80"], fontsize=9.5)
    ax.set_xlabel("mAP bewegter Objekte (AP30–60) · gepoolte "
                  "Testvorhersagen · Achse gezoomt", fontsize=9.3)
    _strip(ax)
    ax.grid(axis="x", color="#EFF0F3", linewidth=1.1)
    ax.set_axisbelow(True)
    fig.tight_layout(pad=0.45)
    fig.savefig(path, dpi=250)
    plt.close(fig)


def fig_dyn_dumbbell(path: Path) -> None:
    """Hantel-Chart je Klasse: AP30 (Ring) -> AP60 (Punkt) je Sicht."""
    fig, axes = plt.subplots(1, 3, figsize=(12.3, 2.9), dpi=250)
    for ax, cls in zip(axes, DYN):
        for row, ((ap30, ap60), color) in enumerate(zip(DYN[cls], VIEW_HEX)):
            ax.plot([ap60, ap30], [row, row], color=CX_HAIR, linewidth=2.4,
                    zorder=1, solid_capstyle="round")
            ax.scatter([ap30], [row], s=95, facecolor="white",
                       edgecolor=color, linewidths=2.0, zorder=3)
            ax.scatter([ap60], [row], s=95, color=color, zorder=3)
            ax.text(ap30 + 0.035, row, f"{ap30:.2f}", va="center", ha="left",
                    fontsize=10, color=CX_SUBTLE)
            ax.text(ap60 - 0.035, row, f"{ap60:.2f}", va="center",
                    ha="right", fontsize=10, fontweight="bold", color=color)
        ax.set_ylim(2.6, -0.6)
        ax.set_xlim(0.05, 1.13)
        ax.set_xticks([0.25, 0.5, 0.75, 1.0])
        ax.set_xticklabels(["0.25", "0.50", "0.75", "1.00"], fontsize=9)
        if ax is axes[0]:
            ax.set_yticks(range(3))
            ax.set_yticklabels(VIEWS, fontsize=11.5)
        else:
            ax.set_yticks([])
        ax.set_title(cls, fontsize=13, fontweight="bold", pad=8,
                     color=CX_INK)
        _strip(ax)
        ax.grid(axis="x", color="#EFF0F3", linewidth=1.1)
        ax.set_axisbelow(True)
    handles = [
        Line2D([], [], marker="o", linestyle="", color=CX_INK, markersize=9,
               label="AP60 — präzise lokalisiert"),
        Line2D([], [], marker="o", linestyle="", markerfacecolor="white",
               markeredgecolor=CX_INK, markeredgewidth=1.8, markersize=9,
               label="AP30 — Objekt gefunden"),
    ]
    fig.legend(handles=handles, loc="upper right", ncol=2, frameon=False,
               fontsize=10, bbox_to_anchor=(0.995, 1.02))
    fig.tight_layout(pad=0.5, w_pad=2.4, rect=(0, 0, 1, 0.93))
    fig.savefig(path, dpi=250)
    plt.close(fig)


def fig_exp_matrix(path: Path) -> None:
    """Punktmatrix: AP30 der bewegten Zielklasse je Experiment × Sicht."""
    fig, ax = plt.subplots(figsize=(12.3, 2.7), dpi=250)
    ramp = LinearSegmentedColormap.from_list(
        "violet", ["#ECE7F9", "#41317F"])
    for row, view in enumerate(("merged", "os0", "os1")):
        for col, value in enumerate(EXP_AP30[view]):
            color = ramp((value - 0.55) / 0.45)
            ax.scatter([col], [row], s=1500, color=color, zorder=2)
            ax.text(col, row, f"{value:.2f}", ha="center", va="center",
                    fontsize=9.5, fontweight="bold", zorder=3,
                    color="white" if value > 0.82 else CX_INK)
    for x in (2.5, 5.5):
        ax.axvline(x, color=CX_HAIR, linewidth=1.2, zorder=1)
    ax.set_xlim(-0.6, 8.6)
    ax.set_ylim(2.7, -0.7)
    ax.set_xticks(range(9))
    ax.set_xticklabels([f"{n}\n{c}" for n, c in EXP_COLS], fontsize=9.5)
    ax.set_yticks(range(3))
    ax.set_yticklabels(VIEWS, fontsize=12)
    _strip(ax, keep_x=False)
    ax.set_xlabel("ungesehenes Testexperiment (bewegte Zielklasse) · "
                  "dunkler = höhere AP30", fontsize=9.5)
    fig.tight_layout(pad=0.4)
    fig.savefig(path, dpi=250)
    plt.close(fig)


def fig_arch_pairs(path: Path) -> None:
    """Paar-Chart: CenterPoint-Punkt -> PointPillars-Punkt je Sicht."""
    fig, ax = plt.subplots(figsize=(6.1, 2.55), dpi=250)
    for row, (view, name) in enumerate(zip(("merged", "os0", "os1"),
                                           VIEWS)):
        cp, pp = CP_MEAN[view], PP_MEAN[view]
        ax.plot([cp, pp], [row, row], color=CX_HAIR, linewidth=2.5,
                zorder=1, solid_capstyle="round")
        ax.scatter([cp], [row], s=120, color=CX_GRAY, zorder=3)
        ax.scatter([pp], [row], s=120, color=CX_ACCENT, zorder=3)
        ax.text(cp - 0.008, row, f"{cp:.3f}", va="center", ha="right",
                fontsize=11, color=CX_SUBTLE)
        ax.text(pp + 0.008, row, f"{pp:.3f}", va="center", ha="left",
                fontsize=11.5, fontweight="bold", color=CX_ACCENT)
        ax.text((cp + pp) / 2, row - 0.32, f"Δ {pp - cp:+.3f}", ha="center",
                fontsize=8.5, color=CX_FAINT)
    ax.set_yticks(range(3))
    ax.set_yticklabels(VIEWS, fontsize=12)
    ax.set_ylim(2.6, -0.75)
    ax.set_xlim(0.60, 0.88)
    ax.set_xticks([0.60, 0.65, 0.70, 0.75, 0.80, 0.85])
    ax.set_xticklabels(["0.60", "0.65", "0.70", "0.75", "0.80", "0.85"],
                       fontsize=9.5)
    ax.set_xlabel("Fold-Mittel mAP · identisches CV-Protokoll · "
                  "Achse gezoomt", fontsize=9.5)
    handles = [
        Line2D([], [], marker="o", linestyle="", color=CX_ACCENT,
               markersize=9, label="PointPillars"),
        Line2D([], [], marker="o", linestyle="", color=CX_GRAY,
               markersize=9, label="CenterPoint"),
    ]
    ax.legend(handles=handles, loc="lower left", frameon=False, fontsize=10,
              ncol=2, bbox_to_anchor=(-0.02, -0.06))
    _strip(ax)
    ax.grid(axis="x", color="#EFF0F3", linewidth=1.1)
    ax.set_axisbelow(True)
    fig.tight_layout(pad=0.4)
    fig.savefig(path, dpi=250)
    plt.close(fig)


def fig_pretrain(path: Path) -> None:
    """Lollipop: Finetuning vs. nur vortrainierte Netze (Vorgängerarbeit)."""
    fig, ax = plt.subplots(figsize=(6.1, 2.35), dpi=250)
    for row, (label, value, ours) in enumerate(PRETRAIN):
        color = CX_ACCENT if ours else CX_GRAY
        ax.plot([0, value], [row, row], color=color, linewidth=2.2,
                zorder=2, solid_capstyle="round")
        ax.scatter([value], [row], s=110, color=color, zorder=3)
        text = "≈ 0" if value == 0 else f"{value:.3f}"
        ax.text(value + 0.02, row, text, va="center", ha="left",
                fontsize=11.5, fontweight="bold",
                color=CX_ACCENT if ours else CX_SUBTLE)
    ax.set_yticks(range(4))
    ax.set_yticklabels([p[0] for p in PRETRAIN], fontsize=10.5)
    ax.set_ylim(3.6, -0.6)
    ax.set_xlim(0, 1.0)
    ax.set_xticks([0, 0.25, 0.5, 0.75, 1.0])
    ax.set_xticklabels(["0", "0.25", "0.50", "0.75", "1.00"], fontsize=9)
    ax.set_xlabel("mAP auf merged · Vorgängerwerte: T. Pagel 2025, "
                  "gleiche Aufnahmen", fontsize=9.5)
    _strip(ax)
    ax.grid(axis="x", color="#EFF0F3", linewidth=1.1)
    ax.set_axisbelow(True)
    fig.tight_layout(pad=0.4)
    fig.savefig(path, dpi=250)
    plt.close(fig)


def fig_runtime(path: Path) -> None:
    """Bullet-Gauge: Inferenzzeit im 100-ms-Budget."""
    fig, ax = plt.subplots(figsize=(6.1, 2.35), dpi=250)
    for row, (view, color, name) in enumerate(zip(("merged", "os0", "os1"),
                                                  VIEW_HEX, VIEWS)):
        ms = RUNTIME_MS[view]
        ax.barh(row, 100, height=0.58, color="#EFF0F3", zorder=1)
        ax.barh(row, ms, height=0.26, color=color, zorder=2)
        ax.text(ms + 2, row, f"{ms} ms", va="center", ha="left",
                fontsize=11.5, fontweight="bold", color=CX_INK)
        ax.text(113, row, f"{ms} %", va="center", ha="right", fontsize=9.5,
                color=CX_FAINT)
    ax.axvline(100, color=CX_INK, linewidth=1.4, zorder=3)
    ax.text(100, -0.72, "Budget: 100 ms (Sensorrate 10 Hz)", ha="right",
            fontsize=9, color=CX_SUBTLE)
    ax.set_yticks(range(3))
    ax.set_yticklabels(VIEWS, fontsize=12)
    ax.set_ylim(2.6, -1.05)
    ax.set_xlim(0, 114)
    ax.set_xticks([0, 25, 50, 75, 100])
    ax.set_xticklabels(["0", "25", "50", "75", "100"], fontsize=9)
    ax.set_xlabel("Inferenzzeit pro Frame in ms (V100, inkl. Voxelisierung) "
                  "· rechts: Budget-Auslastung", fontsize=9.5)
    _strip(ax)
    fig.tight_layout(pad=0.4)
    fig.savefig(path, dpi=250)
    plt.close(fig)


def fig_fp(path: Path) -> None:
    """Gestapelte Geisterboxen aller Klassen je 100 OOF-Frames."""
    fig, ax = plt.subplots(figsize=(5.9, 2.2), dpi=250)
    frames = 2122
    classes = ("Person", "Fahrrad", "Auto")
    for row, (view, name) in enumerate(zip(("merged", "os0", "os1"), VIEWS)):
        left = 0.0
        for class_name, color in zip(classes, GHOST_CLASS_HEX):
            count = GHOST_BOXES[view][class_name]
            rate = 100.0 * count / frames
            ax.barh(row, rate, left=left, height=0.52, color=color,
                    edgecolor="white", linewidth=0.8, zorder=2)
            if rate >= 10:
                ax.text(left + rate / 2, row, f"{rate:.1f}", va="center",
                        ha="center", fontsize=8.7, fontweight="bold",
                        color="white")
            left += rate
        total = sum(GHOST_BOXES[view].values())
        ax.text(left + 1.2, row, f"{left:.1f}  ({total})", va="center",
                ha="left", fontsize=10.2, fontweight="bold", color=CX_INK)
    ax.set_yticks(range(3))
    ax.set_yticklabels(VIEWS, fontsize=12)
    ax.set_ylim(2.6, -0.6)
    ax.set_xlim(0, 85)
    ax.set_xticks([0, 20, 40, 60, 80])
    ax.set_xticklabels(["0", "20", "40", "60", "80"], fontsize=9)
    ax.set_xlabel("Geisterboxen je 100 Frames · rechts: Rate (Gesamtzahl)",
                  fontsize=9.5)
    _strip(ax)
    ax.grid(axis="x", color="#EFF0F3", linewidth=1.1)
    ax.set_axisbelow(True)
    handles = [plt.Rectangle((0, 0), 1, 1, color=color)
               for color in GHOST_CLASS_HEX]
    ax.legend(handles, classes, loc="upper center", bbox_to_anchor=(0.5, -0.29),
              ncol=3, frameon=False, fontsize=9.2, handlelength=1.2,
              columnspacing=1.4)
    fig.tight_layout(pad=0.4)
    fig.savefig(path, dpi=250)
    plt.close(fig)


# ------------------------------------------------------------ Deck-Helfer

def _fill(shape, color, line=None):
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    if line is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line
        shape.line.width = Pt(0.75)
    shape.shadow.inherit = False


def para(frame, first, runs, align=None, after=None, before=None,
         spacing=None):
    p = frame.paragraphs[0] if first else frame.add_paragraph()
    if align is not None:
        p.alignment = align
    if after is not None:
        p.space_after = Pt(after)
    if before is not None:
        p.space_before = Pt(before)
    for text, size, bold, color in runs:
        r = p.add_run()
        r.text = text
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.color.rgb = color
        if spacing is not None:
            r._r.get_or_add_rPr().set("spc", str(spacing))
    return p


def box(slide, x, y, w, h, paras, anchor=None):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = tb.text_frame
    frame.word_wrap = True
    if anchor is not None:
        frame.vertical_anchor = anchor
    for i, spec in enumerate(paras):
        para(frame, i == 0, **spec)
    return tb


def rect(slide, x, y, w, h, color, line=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y),
                                   Inches(w), Inches(h))
    _fill(shape, color, line)
    return shape


def rule(slide, x, y, w):
    rect(slide, x, y, w, 0.012, HAIR)


def kicker(slide, text, x=0.55, y=0.94):
    box(slide, x, y, 8.0, 0.3, [dict(
        runs=[(text.upper(), 10.5, True, FAINT)], spacing=260)])


def content(prs, title):
    slide = prs.slides.add_slide(prs.slide_layouts[8])
    slide.shapes.title.text = title
    for ph in list(slide.placeholders):
        if ph.placeholder_format.idx == 1:
            ph._element.getparent().remove(ph._element)
    return slide


def divider(prs, number, name, sub, image=None, aspect=None):
    slide = content(prs, "")
    box(slide, 0.85, 1.55, 4.0, 2.2, [dict(
        runs=[(number, 150, True, GHOST)])])
    rect(slide, 0.92, 4.15, 0.85, 0.05, ACCENT)
    box(slide, 0.85, 4.42, 7.6, 1.0, [dict(
        runs=[(name, 33, True, INK)])])
    box(slide, 0.85, 5.32, 7.0, 1.0, [dict(
        runs=[(sub, 15, False, SUBTLE)])])
    if image is not None:
        w = 4.9
        slide.shapes.add_picture(str(image), Inches(7.85), Inches(2.25),
                                 width=Inches(w))
    return slide


def check_row(slide, y, requirement, answer):
    box(slide, 0.62, y, 0.4, 0.5, [dict(
        runs=[("✓", 17, True, ACCENT)])])
    box(slide, 1.18, y - 0.05, 6.6, 0.72,
        [dict(runs=[(requirement, 13, True, INK)])],
        anchor=MSO_ANCHOR.MIDDLE)
    box(slide, 7.95, y - 0.05, 4.85, 0.72,
        [dict(runs=[(answer, 11.5, False, SUBTLE)])],
        anchor=MSO_ANCHOR.MIDDLE)


def retain_and_reorder_slides(prs, order):
    """Behält ausgewählte Folien und ordnet Hauptteil vor dem Backup an."""
    id_list = prs.slides._sldIdLst
    slide_ids = list(id_list)
    keep = set(order)
    for index, slide_id in enumerate(slide_ids):
        if index not in keep:
            prs.part.drop_rel(slide_id.rId)
            id_list.remove(slide_id)
    for index in order:
        slide_id = slide_ids[index]
        id_list.remove(slide_id)
        id_list.append(slide_id)


# ---------------------------------------------------------------- Folien

def build(output: Path) -> None:
    tmp_dir = tempfile.TemporaryDirectory(prefix="abschluss_deck_")
    tmp = Path(tmp_dir.name)
    charts = {name: tmp / f"{name}.png" for name in (
        "map_dots", "dynamic_map", "dyn", "exp", "arch", "pretrain",
        "runtime", "fp")}
    fig_map_dots(charts["map_dots"])
    fig_dynamic_map(charts["dynamic_map"])
    fig_dyn_dumbbell(charts["dyn"])
    fig_exp_matrix(charts["exp"])
    fig_arch_pairs(charts["arch"])
    fig_pretrain(charts["pretrain"])
    fig_runtime(charts["runtime"])
    fig_fp(charts["fp"])
    scene_aspect = (lambda im: im.width / im.height)(Image.open(SCENE_MERGED))

    prs = Presentation(TEMPLATE)
    id_list = prs.slides._sldIdLst
    for slide_id in list(id_list):
        prs.part.drop_rel(slide_id.rId)
        id_list.remove(slide_id)

    # ---- 1 · Titel
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.placeholders[0].text = ("Multisensor-LiDAR-Objekterkennung an "
                                  "einer Kreuzung")
    sub = slide.placeholders[1].text_frame
    sub.word_wrap = True
    para(sub, True, [("Lohnt sich die Fusion zweier Sensoren? — "
                      "Abschlusspräsentation des DCAITI-Projekts",
                      15, True, INK)], after=6)
    para(sub, False, [("Juli 2026   ·   Laurenz Gilbert "
                       "(gilbert@campus.tu-berlin.de)   ·   Haoran Wang "
                       "(haoran.wang@campus.tu-berlin.de)",
                       12, False, SUBTLE)])
    box(slide, 1.05, 6.45, 11.2, 0.5, [dict(
        runs=[("Betreuung   ", 12, True, INK),
              ("Johann Nikolai Hark (Fraunhofer FOKUS extern)   ·   "
               "Manuel Schiewe (Fraunhofer FOKUS)", 12, False, SUBTLE)])])

    # ---- 2 · Agenda (vier Spalten)
    slide = content(prs, "Agenda")
    kicker(slide, "Abschlusspräsentation · DCAITI-Projekt")
    cols = [
        ("01", "Projektthema",
         "Aufgabenstellung, Sensor-Setup, Datengrundlage und "
         "Bewertungsprotokoll."),
        ("02", "Projektverlauf",
         "Der ursprüngliche Plan, fünf Wendepunkte und ihre Gründe."),
        ("03", "Stand des Projekts",
         "Cross-Validation-Ergebnisse, Architekturvergleich, Fehlerbild, "
         "Video, Laufzeit, Abgleich mit der Ausschreibung."),
        ("04", "Zusammenfassung & Ausblick",
         "Kernaussagen, offene Punkte und Einstieg für Folgegruppen."),
    ]
    for i, (num, head, body) in enumerate(cols):
        x = 0.55 + i * 3.12
        box(slide, x, 1.65, 2.85, 1.1, [dict(
            runs=[(num, 44, True, GHOST)])])
        rect(slide, x + 0.02, 2.62, 0.55, 0.045, ACCENT)
        box(slide, x, 2.85, 2.85, 0.8, [dict(
            runs=[(head, 15.5, True, INK)])])
        box(slide, x, 3.55, 2.80, 2.4, [dict(
            runs=[(body, 12, False, SUBTLE)])])
    rule(slide, 0.55, 6.35, 12.23)
    box(slide, 0.55, 6.52, 12.2, 0.4, [dict(
        runs=[("Leitfrage durch alle Abschnitte:  ", 12, False, SUBTLE),
              ("Was bringt die Fusion zweier LiDAR-Blickwinkel wirklich?",
               12, True, ACCENT_DEEP)])])

    # ---- 3 · Divider 01
    divider(prs, "01", "Projektthema",
            "Zwei fest montierte LiDARs, eine Kreuzung in einer "
            "Tiefgarage — und die Frage, ob sich das Zusammenführen "
            "ihrer Punktwolken lohnt.", image=SCENE_OS1)

    # ---- 4 · Aufgabenstellung
    slide = content(prs, "Ein Netz soll Person, Fahrrad und Auto in "
                         "Punktwolken finden")
    kicker(slide, "01 · Projektthema — Aufgabe")
    blocks = [
        ("Setup", 1.55,
         "Zwei Ouster-LiDARs (OS0, OS1) beobachten denselben "
         "Kreuzungsbereich aus verschiedenen Richtungen; ihre Punktwolken "
         "werden zusätzlich zu einer fusionierten Sicht (merged) "
         "zusammengeführt."),
        ("Ziel", 2.95,
         "Ein vortrainiertes 3D-Detektionsmodell ausschließlich auf den "
         "eigenen Experimentdaten finetunen und den Nutzen der Fusion "
         "quantifizieren: os0 vs. os1 vs. merged."),
        ("Geforderte Evaluation", 4.35,
         "mAP über die IoU-Schwellen 0.3–0.6 je Klasse, Inferenz- und "
         "Merge-Zeit sowie eine qualitative Analyse von Erfolgen und "
         "Fehlermodi."),
    ]
    for head, y, body in blocks:
        rect(slide, 0.55, y + 0.06, 0.045, 0.92, ACCENT)
        box(slide, 0.75, y, 6.05, 1.4, [
            dict(runs=[(head, 14, True, INK)], after=3),
            dict(runs=[(body, 12.5, False, SUBTLE)]),
        ])
    box(slide, 0.55, 6.05, 6.3, 0.9, [dict(
        runs=[("Kernfrage: ", 14, True, ACCENT_DEEP),
              ("Reicht ein Sensor — oder braucht es beide Blickwinkel?",
               14, True, INK)])])
    img_w = 5.5
    img_h = img_w / scene_aspect
    slide.shapes.add_picture(str(SCENE_MERGED), Inches(7.25), Inches(1.5),
                             width=Inches(img_w))
    box(slide, 7.25, 1.5 + img_h + 0.08, img_w, 0.6, [dict(
        runs=[("Die Kreuzung als fusionierte Punktwolke "
               "(Vogelperspektive) — farbige Boxen: Detektionen des "
               "finetunten Netzes, gestrichelt: Ground Truth",
               10.5, False, FAINT)])])

    # ---- 5 · Datengrundlage
    slide = content(prs, "Ein eigener Datensatz — bewertet auf komplett "
                         "ungesehenen Experimenten")
    kicker(slide, "01 · Projektthema — Daten & Bewertung")
    stats = [("9", "Experimente", "3× Auto · 3× Rad · 3× Person"),
             ("2 122", "Frames je Sicht", "gemeinsame gültige Zeitstempel"),
             ("3", "Sichten", "os0 · os1 · merged (paarweise)"),
             ("≈ 22 : 1", "Auto : Fahrrad", "starkes Klassenungleichgewicht")]
    for i, (value, label, note) in enumerate(stats):
        x = 0.55 + i * 3.12
        box(slide, x, 1.5, 2.9, 0.85, [dict(
            runs=[(value, 37, False, ACCENT_DEEP)])])
        box(slide, x, 2.38, 2.9, 0.32, [dict(
            runs=[(label.upper(), 10, True, FAINT)], spacing=200)])
        box(slide, x, 2.68, 2.9, 0.5, [dict(
            runs=[(note, 11, False, SUBTLE)])])
        if i:
            rect(slide, x - 0.14, 1.62, 0.012, 1.35, HAIR)
    rule(slide, 0.55, 3.45, 12.23)
    box(slide, 0.55, 3.66, 6.4, 0.4, [dict(
        runs=[("3-Fold-Cross-Validation: jedes Experiment genau einmal "
               "komplett im Test", 13, True, INK)])])
    # Fold-Matrix: 3 Folds × 9 Experimente, Testzellen gefüllt
    cell_w, cell_h, gap = 0.56, 0.42, 0.10
    x0, y0 = 1.45, 4.55
    for col, (num, cls) in enumerate(EXP_COLS):
        box(slide, x0 + col * (cell_w + gap) - 0.05, y0 - 0.42,
            cell_w + 0.1, 0.38, [dict(
                runs=[(f"{num} {cls[0]}", 9.5, True, FAINT)],
                align=PP_ALIGN.CENTER)])
    for fold in range(3):
        box(slide, 0.55, y0 + fold * (cell_h + gap) + 0.05, 0.85, 0.35,
            [dict(runs=[(f"Fold {fold + 1}", 10.5, True, SUBTLE)])])
        for col in range(9):
            is_test = col % 3 == fold
            rect(slide, x0 + col * (cell_w + gap),
                 y0 + fold * (cell_h + gap), cell_w, cell_h,
                 ACCENT if is_test else TRACK)
    lx = x0 + 9 * (cell_w + gap) + 0.25
    rect(slide, lx, y0 + 0.02, 0.26, 0.26, ACCENT)
    box(slide, lx + 0.36, y0 - 0.03, 1.6, 0.35, [dict(
        runs=[("Test (ungesehen)", 10, False, SUBTLE)])])
    rect(slide, lx, y0 + 0.54, 0.26, 0.26, TRACK)
    box(slide, lx + 0.36, y0 + 0.49, 1.6, 0.35, [dict(
        runs=[("Training + Val", 10, False, SUBTLE)])])
    box(slide, 0.55, 6.25, 12.2, 0.9, [
        dict(runs=[("Labels aller drei Sichten manuell korrigiert "
                    "(Klassen, Boxmaße, bewegt/statisch) · bewertet werden "
                    "die 2 122 Zeitstempel, die in allen Sichten gültig "
                    "sind — jede Sicht sieht dieselben Szenen.",
                    12, False, SUBTLE)])])

    # ---- 6 · Divider 02
    divider(prs, "02", "Projektverlauf",
            "Die geplante Pipeline trug von Anfang bis Ende — aber fünf "
            "Dinge kamen anders als gedacht.")

    # ---- 7 · Wendepunkte
    slide = content(prs, "Die Pipeline stand früh — die Qualität musste "
                         "erkämpft werden")
    kicker(slide, "02 · Projektverlauf")
    box(slide, 0.55, 1.42, 12.2, 0.65, [
        dict(runs=[("Plan:  ", 12.5, True, INK),
                   ("ROS2-Export  →  Punktwolken-Merge (ICP)  →  "
                    "manuelles Labeling  →  KITTI-Konvertierung  →  "
                    "PointPillars-Finetuning  →  Evaluation in drei "
                    "Sichten", 12.5, False, SUBTLE)], after=2),
        dict(runs=[("Alle sechs Schritte wurden umgesetzt — die "
                    "Wendepunkte lagen in Daten- und Messqualität:",
                    12.5, True, ACCENT_DEEP)]),
    ])
    turns = [
        ("Automatische Vor-Labels zu ungenau",
         "alle 9 Experimente in allen 3 Sichten von Hand nachgezogen — "
         "inkl. bewegt/statisch-Kennzeichnung."),
        ("Drei versteckte Datenfehler im Audit gefunden",
         "abgeschnittener Boden, schwebende Boxen, inkonsistente "
         "Intensitäten → Konverter gefixt, Daten v2, Neutraining."),
        ("Übernommene Messmethode verzerrte seltene Klassen",
         "Frame-Mittelung → datensatzweite Standard-AP neu implementiert."),
        ("Klassenungleichgewicht 22:1",
         "GT-Sampling übernommen; Alternativen getestet und als "
         "Negativergebnisse dokumentiert."),
        ("Temporaler Test-Split mit Objektposition konfundiert",
         "Bewertung auf experiment-basierte Cross-Validation umgestellt — "
         "erst sie misst echte Generalisierung."),
    ]
    line_x = 0.85
    y_first, step = 2.62, 0.88
    rect(slide, line_x - 0.008, y_first + 0.16, 0.016,
         step * 4 + 0.05, HAIR)
    for i, (head, fix) in enumerate(turns):
        y = y_first + i * step
        shape = slide.shapes.add_shape(
            MSO_SHAPE.OVAL, Inches(line_x - 0.17), Inches(y),
            Inches(0.34), Inches(0.34))
        _fill(shape, ACCENT)
        frame = shape.text_frame
        frame.word_wrap = False
        para(frame, True, [(str(i + 1), 12, True, WHITE)],
             align=PP_ALIGN.CENTER)
        box(slide, 1.30, y - 0.10, 11.4, 0.85, [
            dict(runs=[(head + "   ", 13, True, INK),
                       ("— " + fix, 12, False, SUBTLE)]),
        ])

    # ---- 8 · Messung
    slide = content(prs, "Der teuerste Fehler steckte nicht im Netz, "
                         "sondern in der Messung")
    kicker(slide, "Backup · Daten- und Messqualität")
    rect(slide, 0.55, 1.55, 5.7, 2.75, TINT)
    box(slide, 0.85, 1.85, 5.2, 0.35, [dict(
        runs=[("FAHRRAD-AP30 · DERSELBE MODELLOUTPUT", 10, True,
               ACCENT_DEEP)], spacing=200)])
    box(slide, 0.85, 2.25, 5.2, 1.1, [dict(
        runs=[("0.24", 40, True, FAINT), ("  →  ", 28, False, SUBTLE),
              ("0.94", 40, True, ACCENT_DEEP)])])
    box(slide, 0.85, 3.35, 5.1, 0.8, [dict(
        runs=[("vorher: AP je Frame gemittelt — Frames ohne Fahrrad "
               "zählten als 0 · nachher: datensatzweite Standard-AP",
               11.5, False, SUBTLE)])])
    box(slide, 0.55, 4.65, 5.7, 1.6, [
        dict(runs=[("Aufgefallen, weil zwei verschieden trainierte Modelle "
                    "exakt denselben Fahrrad-Wert erreichten — genau das "
                    "Maximum der fehlerhaften Rechnung.", 12, False,
                    SUBTLE)], after=6),
        dict(runs=[("Der Fehler machte das Netz nie schlechter — er hat "
                    "nur verdeckt, wie gut es war.", 12, True, INK)]),
    ])
    box(slide, 6.85, 1.5, 5.9, 0.4, [dict(
        runs=[("Drei Datenfehler, im Audit gefunden und behoben",
               13.5, True, INK)])])
    audit = [
        ("Boden abgeschnitten", 2.05,
         "Ein Höhenversatz schob die Szene unter den Erfassungsbereich — "
         "Boden und untere 40 cm aller Objekte fehlten im Training."),
        ("Boxen schwebten", 3.20,
         "Labelhöhe war als Boxmitte gespeichert, das Framework las sie "
         "als Unterkante — alle Boxen saßen ~75 cm zu hoch."),
        ("Intensitäten nicht vergleichbar", 4.35,
         "Normierung pro Frame statt global — dieselbe Oberfläche bekam "
         "je Aufnahme andere Werte."),
    ]
    for head, y, body in audit:
        rect(slide, 6.85, y + 0.06, 0.045, 0.8, ACCENT)
        box(slide, 7.05, y, 5.75, 1.15, [
            dict(runs=[(head, 13, True, INK)], after=2),
            dict(runs=[(body, 11.5, False, SUBTLE)]),
        ])
    box(slide, 6.85, 5.6, 5.9, 0.7, [dict(
        runs=[("Konsequenz: ", 12.5, True, ACCENT_DEEP),
              ("alle Ergebnisse vor den Fixes verworfen — jede gezeigte "
               "Zahl basiert auf dem bereinigten Datenstand v2.",
               12.5, False, INK)])])

    # ---- 9 · Divider 03
    divider(prs, "03", "Stand des Projekts",
            "Alle Anforderungen der Ausschreibung sind umgesetzt — "
            "evaluiert per experiment-basierter Cross-Validation, "
            "zusätzlich mit zweiter Architektur gegengeprüft.")

    # ---- 10 · Gesamtergebnis
    slide = content(prs, "Die Fusion liefert in jedem Fold den besten "
                         "Gesamtwert — knapp, aber konsistent")
    kicker(slide, "03 · Projektstand — Gesamtergebnis")
    box(slide, 0.55, 1.5, 6.2, 0.35, [dict(
        runs=[("Gesamt-mAP je Sicht, 3-Fold-Cross-Validation", 12.5, True,
               INK)])])
    slide.shapes.add_picture(str(charts["map_dots"]), Inches(0.55),
                             Inches(1.95), width=Inches(6.2))
    reads = [
        ("In allen drei Folds vorn", 1.95,
         "merged 0.798 · os0 0.777 · os1 0.770 (Fold-Mittel) — der "
         "Vorsprung ist moderat, aber er dreht sich nie um."),
        ("Alle Sichten brauchbar", 3.15,
         "Auch die Einzelsensoren liefern ein solides Gesamtniveau — "
         "kein Sensor fällt aus."),
        ("Gesamtwert allein täuscht", 4.35,
         "Die Zahlen werden von wiedererkannten geparkten Autos getragen "
         "— entscheidend sind die bewegten Objekte (nächste Folie)."),
    ]
    for head, y, body in reads:
        rect(slide, 7.15, y + 0.06, 0.045, 0.85, ACCENT)
        box(slide, 7.35, y, 5.45, 1.2, [
            dict(runs=[(head, 13.5, True, INK)], after=2),
            dict(runs=[(body, 12, False, SUBTLE)]),
        ])
    rule(slide, 0.55, 6.0, 12.23)
    box(slide, 0.55, 6.18, 12.2, 0.5, [dict(
        runs=[("Bewertung ausschließlich auf Experimenten, die das "
               "jeweilige Modell nie gesehen hat.", 12, False, FAINT)])])

    # ---- 10b · Dynamischer Gesamtwert
    slide = content(prs, "Auch bei bewegten Objekten liefert die Fusion "
                         "den höchsten Gesamtwert")
    kicker(slide, "03 · Projektstand — dynamischer Gesamt-mAP")
    box(slide, 0.55, 1.5, 6.1, 0.5, [dict(
        runs=[("mAP nur für bewegte Objekte, alle drei Klassen",
               12.5, True, INK)])])
    slide.shapes.add_picture(str(charts["dynamic_map"]), Inches(0.55),
                             Inches(1.95), width=Inches(6.1))
    dynamic_reads = [
        ("Fusion führt mit 0.836", 1.85,
         "os1 erreicht 0.756, os0 0.710 — der Abstand ist bei den "
         "bewegten Objekten deutlicher als im Gesamtwert aller Objekte."),
        ("Einzelsensoren sind spezialisiert", 3.15,
         "os0 ist beim bewegten Auto stark, verliert aber bei Personen. "
         "Die Fusion bleibt über Person, Fahrrad und Auto gleichmäßig."),
        ("Zusammenfassung, kein neuer Fold-Wert", 4.65,
         "Die Vorhersagen der drei ungesehenen Testfolds sind hier "
         "gepoolt. Die Fold-Ergebnisse der vorherigen Folie bleiben die "
         "primäre Generalisierungsschätzung."),
    ]
    for head, y, body in dynamic_reads:
        rect(slide, 7.15, y + 0.06, 0.045, 0.9, ACCENT)
        box(slide, 7.35, y, 5.45, 1.35, [
            dict(runs=[(head, 13.5, True, INK)], after=2),
            dict(runs=[(body, 12, False, SUBTLE)]),
        ])
    rule(slide, 0.55, 6.25, 12.23)
    box(slide, 0.55, 6.42, 12.2, 0.5, [dict(
        runs=[("mAP = Mittel über AP30/40/50/60 und die drei bewegten "
               "Klassen · Bewegt-Wertung nutzt die bewegt/statisch-Labels: "
               "Boxen auf statisch gelabelten Objekten zählen nicht als "
               "Fehlalarm · Auto ist mit n = 339 exakt gepaart.",
               10.8, False, FAINT)])])

    # ---- 11 · Bewegte Objekte
    slide = content(prs, "Gefunden wird überall — die Präzision "
                         "entscheidet den Vergleich")
    kicker(slide, "03 · Projektstand — bewegte Objekte in ungesehenen Tests")
    slide.shapes.add_picture(str(charts["dyn"]), Inches(0.55), Inches(1.55),
                             width=Inches(12.3))
    box(slide, 0.55, 4.85, 12.2, 1.7, [
        dict(runs=[("Detektion (Ringe): ", 12.5, True, INK),
                   ("jede Sicht findet jede bewegte Klasse — AP30 "
                    "durchgehend ≥ 0.65, meist ≥ 0.80. Auch das bewegte "
                    "Auto, in jeder Sicht.", 12.5, False, SUBTLE)],
             after=6),
        dict(runs=[("Lokalisierung (Punkte): ", 12.5, True, INK),
                   ("bei AP60 liegt die Fusion in allen drei Klassen vorn "
                    "— deutlich bei Person (0.73 vs. 0.37/0.43) und "
                    "Fahrrad (0.67), knapp beim Auto (0.88 vs. 0.86).",
                    12.5, False, SUBTLE)], after=6),
        dict(runs=[("Der belastbare Fusionsvorteil ist die Präzision der "
                    "Lokalisierung.", 12.5, True, ACCENT_DEEP)]),
    ])
    box(slide, 0.55, 6.75, 12.2, 0.4, [dict(
        runs=[("Bewegt-Wertung nutzt die bewegt/statisch-Labels: Boxen "
               "auf statisch gelabelten Objekten zählen nicht als "
               "Fehlalarm.", 10.8, False, FAINT)])])

    # ---- 12 · Robustheit
    slide = content(prs, "27 Zielklassen-Bewertungen auf ungesehenen "
                         "Experimenten — kein Ausfall")
    kicker(slide, "03 · Projektstand — Generalisierung je Experiment")
    slide.shapes.add_picture(str(charts["exp"]), Inches(0.55), Inches(1.55),
                             width=Inches(12.3))
    box(slide, 0.55, 4.75, 12.2, 1.5, [
        dict(runs=[("Schlechtester Wert im gesamten Raster: 0.72 ", 12.5,
                    True, INK),
                   ("(os1, Rad-Experiment 5). Der frühere Ausreißer "
                    "„os1 erkennt das Auto nicht (0.09)“ war ein Artefakt "
                    "des alten temporalen Splits — unter fairer Bewertung "
                    "existiert er nicht.", 12.5, False, SUBTLE)], after=6),
        dict(runs=[("Die Fusion ist in keiner Spalte die schwächste Sicht "
                    "und bleibt auch in den schwierigen Rad-Experimenten "
                    "5 und 6 stabil.", 12.5, False, SUBTLE)]),
    ])

    # ---- 13 · Architekturvergleich
    slide = content(prs, "Zweite Architektur, gleicher Befund — "
                         "PointPillars bleibt vorn")
    kicker(slide, "03 · Projektstand — Architekturvergleich")
    box(slide, 0.55, 1.5, 6.1, 0.35, [dict(
        runs=[("Fold-Mittel mAP: CenterPoint → PointPillars, je Sicht",
               12.5, True, INK)])])
    slide.shapes.add_picture(str(charts["arch"]), Inches(0.55),
                             Inches(1.95), width=Inches(6.1))
    box(slide, 0.55, 4.85, 6.1, 1.4, [dict(
        runs=[("Nebenbefund: ", 12, True, INK),
              ("CenterPoint produziert weniger Geister-Fahrräder "
               "(77–239 vs. 157–372 je Sicht), gleicht den "
               "Präzisionsrückstand damit aber nicht aus.", 12, False,
               SUBTLE)])])
    arch_reads = [
        ("18 Trainingsläufe, ein Protokoll", 1.95,
         "CenterPoint (anchor-frei, Center-Heatmap) unter exakt demselben "
         "CV-Protokoll trainiert wie PointPillars (anchor-basiert)."),
        ("PointPillars in jeder Sicht vorn", 3.15,
         "Δ 0.05–0.12 mAP; Haupttreiber ist die durchweg schwächere "
         "AP60-Lokalisierung — am deutlichsten beim Fahrrad."),
        ("Fusionsbefund architekturunabhängig", 4.35,
         "Auch bei CenterPoint ist merged in jedem Fold die beste Sicht — "
         "und der scheinbare Auto-Kollaps im alten Split trat bei beiden "
         "Architekturen auf und verschwand bei beiden unter der CV."),
    ]
    for head, y, body in arch_reads:
        rect(slide, 7.15, y + 0.06, 0.045, 0.85, ACCENT)
        box(slide, 7.35, y, 5.45, 1.35, [
            dict(runs=[(head, 13.5, True, INK)], after=2),
            dict(runs=[(body, 12, False, SUBTLE)]),
        ])

    # ---- 14 · Fehlerbild
    slide = content(prs, "Geisterboxen zeigen die Grenze der "
                         "Ein-Szenen-Daten")
    kicker(slide, "03 · Projektstand — Fehlerbild")
    box(slide, 0.55, 1.5, 5.9, 0.35, [dict(
        runs=[("Geisterboxen aller Klassen in ungesehenen Tests",
               12.5, True, INK)])])
    slide.shapes.add_picture(str(charts["fp"]), Inches(0.55), Inches(1.95),
                             width=Inches(5.9))
    box(slide, 0.55, 4.55, 5.9, 1.7, [dict(
        runs=[("Gezählt wird jede Vorhersage ab Konfidenz 0,3 ohne passende "
               "Ground-Truth-Box derselben Klasse bei IoU ≥ 0,3. Damit "
               "enthält das Schaubild Personen-, Fahrrad- und Auto-Boxen — "
               "nicht nur die auffälligen Fahrräder.",
               12, False, SUBTLE)])])
    err = [
        ("Alle drei Klassen betroffen", 1.5,
         "merged: 984 · os0: 1.520 · os1: 949 Geisterboxen auf jeweils "
         "2.122 identischen Testframes. Bei os0 dominieren Personen- und "
         "Auto-Boxen, bei os1 fallen besonders Fahrräder auf."),
        ("Mehrere beobachtete Ursachen", 3.0,
         "nicht gelabelte reale Strukturen · ungenaue oder doppelte Boxen · "
         "Reflexions- und Bodencluster · Positionsmemorierung an bekannten "
         "Routen. Nicht jede Geisterbox ist eine leere Halluzination."),
        ("Praktisch unterdrückbar", 4.5,
         "Viele Boxen haben niedrige Konfidenzen oder bestehen nur einen "
         "Frame. Tracking, eine Hintergrundkarte und klassenübergreifende "
         "Bereinigung können die Ausgabe deutlich beruhigen."),
    ]
    for head, y, body in err:
        rect(slide, 6.95, y + 0.06, 0.045, 1.0, ACCENT)
        box(slide, 7.15, y, 5.65, 1.5, [
            dict(runs=[(head, 13.5, True, INK)], after=2),
            dict(runs=[(body, 12, False, SUBTLE)]),
        ])

    # ---- 15 · Video
    slide = content(prs, "Drei Sichten — Inferenz auf demselben Experiment")
    kicker(slide, "03 · Projektstand — qualitativ")
    vid_w, vid_h = 11.0, 11.0 * 9 / 16
    slide.shapes.add_movie(str(VIDEO), Inches((13.33 - vid_w) / 2),
                           Inches(1.28), Inches(vid_w), Inches(vid_h),
                           poster_frame_image=str(POSTER),
                           mime_type="video/mp4")

    # ---- 16 · Einordnung + Laufzeit
    slide = content(prs, "Finetuning macht den Unterschied — die "
                         "Modellinferenz passt in den Sensortakt")
    kicker(slide, "03 · Projektstand — Einordnung & Laufzeit")
    box(slide, 0.55, 1.5, 6.1, 0.6, [
        dict(runs=[("Gegen die Vorgängerarbeit: Faktor ", 13, True, INK),
                   ("≈ 14", 13, True, ACCENT_DEEP)], after=1),
        dict(runs=[("gleiche Aufnahmen, dort ohne Finetuning "
                    "(Größenordnungsvergleich)", 11, False, FAINT)]),
    ])
    slide.shapes.add_picture(str(charts["pretrain"]), Inches(0.55),
                             Inches(2.25), width=Inches(6.1))
    box(slide, 6.95, 1.5, 5.85, 0.6, [
        dict(runs=[("Inferenz im 100-ms-Takt der Sensoren", 13, True,
                    INK)], after=1),
        dict(runs=[("V100, Batch 1, inkl. Datenladen und Voxelisierung",
                    11, False, FAINT)]),
    ])
    slide.shapes.add_picture(str(charts["runtime"]), Inches(6.95),
                             Inches(2.25), width=Inches(5.85))
    rule(slide, 0.55, 5.35, 12.23)
    box(slide, 0.55, 5.55, 12.2, 1.1, [
        dict(runs=[("Ein fertig vortrainiertes Netz direkt anzuwenden "
                    "scheitert an dieser Szene — das Nachtraining auf den "
                    "Zieldaten ist der entscheidende Schritt.", 12.5, False,
                    SUBTLE)], after=5),
        dict(runs=[("Der Offline-Merge benötigt 2,25 s/Frame. Mit fest "
                    "kalibrierter Extrinsik könnte ICP durch eine reine "
                    "Transformation ersetzt werden; diese Online-Variante "
                    "und die End-to-end-Laufzeit sind noch zu messen.",
                    12.5, False, SUBTLE)]),
    ])

    # ---- 17 · Ausschreibung
    slide = content(prs, "Alle Punkte der Projektausschreibung sind "
                         "umgesetzt")
    kicker(slide, "03 · Projektstand — Abgleich mit der Ausschreibung")
    rows = [
        ("Vortrainiertes Modell, Finetuning nur auf Experimentdaten",
         "PointPillars ab KITTI-Checkpoint · 50 Epochen je Sicht und Fold"),
        ("Einzelsensor vs. merged als getrennte Eingaben",
         "os0 / os1 / merged mit identischen Zeitstempeln, je 3 "
         "CV-Trainings"),
        ("mAP über die IoU-Schwellen 0.3 / 0.4 / 0.5 / 0.6 je Klasse",
         "vollständige Tabellen: results/CROSS_VALIDATION_RESULTS.md"),
        ("Inferenzzeit und Merge-Zeit",
         "43–85 ms Modellinferenz · Merge offline 2,25 s; online noch offen"),
        ("Qualitative Analyse von Erfolgen und Fehlermodi",
         "BEV-Analysen, Video, drei bis zur Ursache verfolgte Fehlermodi"),
        ("Optionale Vorverarbeitung (Rauschen, Deckenentfernung)",
         "Rauschfilter aktiv · Decke bewusst behalten — das finetunte "
         "Netz lernt sie als Hintergrund"),
    ]
    y = 1.55
    for requirement, answer in rows:
        check_row(slide, y, requirement, answer)
        rule(slide, 0.55, y + 0.68, 12.23)
        y += 0.84
    box(slide, 0.55, y + 0.05, 12.2, 0.45, [dict(
        runs=[("Darüber hinaus: ", 12, True, ACCENT_DEEP),
              ("Daten-Audit, Metrik-Korrektur, Ablationen, "
               "Cross-Validation-Protokoll und Architekturvergleich.",
               12, False, SUBTLE)])])

    # ---- 18 · Divider 04
    divider(prs, "04", "Zusammenfassung & Ausblick",
            "Drei belegte Kernaussagen — und die Anschlussstellen, an "
            "denen Folgegruppen direkt weiterarbeiten können.")

    # ---- 19 · Kernaussagen
    slide = content(prs, "Drei Aussagen, die dieses Projekt belegt")
    kicker(slide, "04 · Zusammenfassung")
    findings = [
        ("Ohne Finetuning geht es nicht",
         "Vortrainierte Netze erreichen auf dieser Szene mAP ≈ 0.06 — "
         "nach dem Finetuning 0.80. Der Schritt auf die Zieldaten ist "
         "der Hebel, nicht die Architektur.", "≈ 14×"),
        ("Fusion ist die beste Balance — nicht die Rettung",
         "Bester Gesamt-mAP in jedem Fold beider Architekturen, nie die "
         "schwächste Sicht, präziseste Lokalisierung bewegter Objekte in "
         "allen drei Klassen. Das bewegte Auto findet aber jede Sicht.",
         "0.80"),
        ("Die Modellinferenz ist 10-Hz-fähig",
         "Die fusionierte Punktwolke bleibt mit 85 ms unter dem "
         "100-ms-Takt. Für ein Live-Gesamtsystem müssen Online-Merge und "
         "End-to-end-Laufzeit noch validiert werden.", "85 ms"),
    ]
    y = 1.6
    for i, (head, body, metric) in enumerate(findings):
        box(slide, 0.55, y - 0.12, 1.1, 1.4, [dict(
            runs=[(str(i + 1), 54, True, GHOST)])])
        box(slide, 1.75, y, 8.3, 1.5, [
            dict(runs=[(head, 16, True, INK)], after=4),
            dict(runs=[(body, 12.5, False, SUBTLE)]),
        ])
        box(slide, 10.3, y - 0.05, 2.45, 1.0,
            [dict(runs=[(metric, 30, True, ACCENT_DEEP)],
                  align=PP_ALIGN.RIGHT)])
        if i < 2:
            rule(slide, 0.55, y + 1.42, 12.23)
        y += 1.72


    # ---- 20 · Ausblick
    slide = content(prs, "Ausblick: fünf konkrete nächste Schritte")
    kicker(slide, "04 · Ausblick")
    box(slide, 0.55, 1.40, 12.2, 0.5, [dict(
        runs=[("Von der wissenschaftlichen Bewertung zum einsatzfähigen "
               "und besser generalisierenden System", 12.5, False,
               SUBTLE)])])
    outlook_items = [
        ("Finales Einsatzmodell",
         "gewählte Sicht auf allen gültigen Daten neu trainieren; die "
         "Cross-Validation bleibt die Leistungsschätzung."),
        ("Mehr Szenen und mehr Objekte",
         "die Ein-Szenen-Datenbasis erweitern und Positions-Memorierung "
         "durch echte Szenenvielfalt reduzieren."),
        ("Tracking über mehrere Frames",
         "Objekte zeitlich verbinden, kurzzeitige Aussetzer überbrücken "
         "und Ein-Frame-Fehldetektionen filtern."),
        ("Online-Merge",
         "einmal kalibrierte Extrinsik statt ICP pro Frame verwenden und "
         "die End-to-end-Laufzeit messen."),
        ("Hintergrundkarte der Szene",
         "dauerhaft bekannte Strukturen nutzen, um ortsgebundene "
         "Fehldetektionen zu unterdrücken."),
    ]
    y = 1.92
    for index, (head, body) in enumerate(outlook_items, start=1):
        box(slide, 0.55, y - 0.08, 0.55, 0.65, [dict(
            runs=[(str(index), 28, True, GHOST)],
            align=PP_ALIGN.CENTER)])
        rect(slide, 1.25, y + 0.05, 0.045, 0.55, ACCENT)
        box(slide, 1.48, y, 11.2, 0.72, [dict(
            runs=[(head + " — ", 12.8, True, INK),
                  (body, 12.2, False, SUBTLE)])])
        if index < len(outlook_items):
            rule(slide, 1.48, y + 0.73, 11.2)
        y += 0.88
    rule(slide, 0.55, 6.35, 12.23)
    box(slide, 0.55, 6.52, 12.2, 0.45, [dict(
        runs=[("Alles reproduzierbar: Configs, Splits und Auswertungen im "
               "Repo · CV-Checkpoints beider Architekturen auf dem "
               "Projektserver · ältere Modelle als GitHub-Release v1.0.",
               11.5, False, FAINT)])])

    # ---- 21 · Danke
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.placeholders[0].text = "Vielen Dank — Fragen?"
    sub = slide.placeholders[1].text_frame
    sub.word_wrap = True
    para(sub, True, [("Laurenz Gilbert · gilbert@campus.tu-berlin.de      "
                      "Haoran Wang · haoran.wang@campus.tu-berlin.de",
                      12, False, SUBTLE)])

    box(slide, 1.05, 6.65, 11.2, 0.35, [dict(
        runs=[("Backupfolien ab der nächsten Folie", 10.5, False, FAINT)],
        align=PP_ALIGN.RIGHT)])

    # 16 kompakte Inhaltsfolien + 4 kurze Abschnittstrenner + Fragenfolie
    # + 1 Backupfolie. Die Trenner werden im Vortrag nur kurz eingeblendet.
    # Die Indizes beziehen sich auf die oben erzeugte 22-Folien-Fassung.
    main_slides = [0, 1, 2, 3, 4, 5, 6, 8, 17, 9, 10, 11, 12, 13, 14, 15,
                   16, 18, 19, 20, 21]
    backup_slides = [7]
    retain_and_reorder_slides(prs, main_slides + backup_slides)

    prs.save(output)
    tmp_dir.cleanup()


def main() -> None:
    parser = argparse.ArgumentParser(description=__doc__)
    parser.add_argument("--output", type=Path, default=DEFAULT_OUTPUT)
    args = parser.parse_args()
    build(args.output)
    print(f"Deck geschrieben: {args.output}")


if __name__ == "__main__":
    main()
